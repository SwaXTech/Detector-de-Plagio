from pptx import Presentation as ppt
from util.data_cleaning import separate_glued_words, remove_multiple_whitespaces, merge_string

class Presentation():

  presentation = None
  string = None

  def __init__(self, path):
    self.presentation = ppt(path)
    self.string = self.build_string()

  def raw_full_text(self):
    raw_text = []
    for slide in self.presentation.slides:
      for shape in slide.shapes:
          if not shape.has_text_frame:
              continue
          for paragraph in shape.text_frame.paragraphs:
              for child in paragraph.runs: ## Text run object. Corresponds to a child element in a paragraph.
                  raw_text.append(child.text)

    return raw_text

  def trim_text(self, raw_text):
    raw_text = separate_glued_words(raw_text)
    return remove_multiple_whitespaces(raw_text)

  def build_string(self):
    return self.trim_text(
        merge_string(
            self.raw_full_text()
        )
    )
